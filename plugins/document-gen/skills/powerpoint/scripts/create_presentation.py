#!/usr/bin/env python3
"""
PowerPoint Presentation Generator

Template script for creating professional presentations using python-pptx.
Customize the STYLE and SLIDES sections to match your requirements.

Usage:
    uv run python create_presentation.py [output_path]

Examples:
    uv run python create_presentation.py
    uv run python create_presentation.py ~/Desktop/my_deck.pptx
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


# =============================================================================
# STYLE CONFIGURATION - Customize colors and fonts here
# =============================================================================

class Style:
    """Presentation style settings."""

    # Primary colors - use bold, saturated colors (avoid pastels)
    PRIMARY = RGBColor(0, 82, 147)         # #005293 - Strong blue
    PRIMARY_DARK = RGBColor(0, 62, 117)    # Darker variant
    SECONDARY = RGBColor(214, 102, 75)     # #D6664B - Coral accent
    ACCENT = RGBColor(218, 175, 65)        # #DAAF41 - Gold highlight

    # Backgrounds
    BG_WHITE = RGBColor(255, 255, 255)     # Pure white
    BG_LIGHT = RGBColor(248, 249, 250)     # #F8F9FA - Light gray
    BG_DARK = RGBColor(33, 37, 41)         # #212529 - Near black

    # Text colors
    TEXT_PRIMARY = RGBColor(33, 37, 41)    # #212529 - Dark text
    TEXT_SECONDARY = RGBColor(108, 117, 125)  # #6C757D - Gray text
    TEXT_ON_DARK = RGBColor(255, 255, 255) # White text for dark backgrounds

    # Typography - use cross-platform safe fonts
    TITLE_FONT = "Arial"
    BODY_FONT = "Arial"
    TITLE_SIZE = Pt(44)
    SUBTITLE_SIZE = Pt(22)
    HEADING_SIZE = Pt(34)
    BODY_SIZE = Pt(20)

    # Slide dimensions (16:9 widescreen)
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)

    # Layout
    MARGIN_LEFT = Inches(0.8)
    CONTENT_WIDTH = Inches(11.7)


# =============================================================================
# CONTENT CONFIGURATION - Define your slides here
# =============================================================================

PRESENTATION_TITLE = "Presentation Title"
PRESENTATION_SUBTITLE = "Your Name | Date"
OUTPUT_FILENAME = "presentation.pptx"

# Supported slide types: "title", "section", "content", "closing", "blank"
SLIDES = [
    {
        "type": "title",
        "title": PRESENTATION_TITLE,
        "subtitle": PRESENTATION_SUBTITLE,
    },
    {
        "type": "section",
        "title": "Section 1",
    },
    {
        "type": "content",
        "title": "Key Points",
        "bullets": [
            "First important point",
            "Second important point",
            "Third important point",
        ],
    },
    {
        "type": "content",
        "title": "Details",
        "bullets": [
            "Supporting information",
            "Additional context",
            "Final thoughts",
        ],
    },
    {
        "type": "section",
        "title": "Conclusion",
    },
    {
        "type": "content",
        "title": "Summary",
        "bullets": [
            "Key takeaway 1",
            "Key takeaway 2",
            "Next steps and call to action",
        ],
    },
    {
        "type": "closing",
        "title": "Thank You",
        "subtitle": "contact@example.com",
    },
]


# =============================================================================
# SLIDE BUILDERS
# =============================================================================

def create_presentation() -> Presentation:
    """Create and return a styled presentation."""
    prs = Presentation()
    prs.slide_width = Style.SLIDE_WIDTH
    prs.slide_height = Style.SLIDE_HEIGHT

    for config in SLIDES:
        slide_type = config.get("type", "content")

        if slide_type == "title":
            add_title_slide(prs, config)
        elif slide_type == "section":
            add_section_slide(prs, config)
        elif slide_type == "content":
            add_content_slide(prs, config)
        elif slide_type == "closing":
            add_closing_slide(prs, config)
        elif slide_type == "blank":
            add_blank_slide(prs, config)

    return prs


def add_background(slide, color: RGBColor):
    """Set slide background to a solid color."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs: Presentation, config: dict):
    """Title slide with vertical accent bar."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_background(slide, Style.BG_WHITE)

    # Vertical accent bar on left edge
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Inches(0), top=Inches(0),
        width=Inches(0.35), height=Style.SLIDE_HEIGHT
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = Style.PRIMARY
    bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1.2), Inches(2.6), Inches(10.5), Inches(1.2)
    )
    p = title_box.text_frame.paragraphs[0]
    p.text = config.get("title", "Title")
    p.font.size = Style.TITLE_SIZE
    p.font.bold = True
    p.font.color.rgb = Style.TEXT_PRIMARY
    p.font.name = Style.TITLE_FONT

    # Subtitle
    if "subtitle" in config:
        sub_box = slide.shapes.add_textbox(
            Inches(1.2), Inches(4.0), Inches(10.5), Inches(0.6)
        )
        p = sub_box.text_frame.paragraphs[0]
        p.text = config["subtitle"]
        p.font.size = Style.SUBTITLE_SIZE
        p.font.color.rgb = Style.TEXT_SECONDARY
        p.font.name = Style.BODY_FONT

    # Horizontal accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Inches(1.2), top=Inches(4.8),
        width=Inches(2.5), height=Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = Style.SECONDARY
    line.line.fill.background()


def add_section_slide(prs: Presentation, config: dict):
    """Section divider with solid primary background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, Style.PRIMARY)

    # Section title - white on primary background
    title_box = slide.shapes.add_textbox(
        Style.MARGIN_LEFT, Inches(3.0), Style.CONTENT_WIDTH, Inches(1.5)
    )
    p = title_box.text_frame.paragraphs[0]
    p.text = config.get("title", "Section")
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = Style.TEXT_ON_DARK
    p.font.name = Style.TITLE_FONT

    # Accent line below title
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Style.MARGIN_LEFT, top=Inches(4.6),
        width=Inches(1.5), height=Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = Style.ACCENT
    line.line.fill.background()


def add_content_slide(prs: Presentation, config: dict):
    """Content slide with top accent line and square bullet markers."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, Style.BG_WHITE)

    # Top border accent line
    top_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Inches(0), top=Inches(0),
        width=Style.SLIDE_WIDTH, height=Inches(0.04)
    )
    top_line.fill.solid()
    top_line.fill.fore_color.rgb = Style.PRIMARY
    top_line.line.fill.background()

    # Slide title
    title_box = slide.shapes.add_textbox(
        Style.MARGIN_LEFT, Inches(0.6), Style.CONTENT_WIDTH, Inches(0.9)
    )
    p = title_box.text_frame.paragraphs[0]
    p.text = config.get("title", "")
    p.font.size = Style.HEADING_SIZE
    p.font.bold = True
    p.font.color.rgb = Style.TEXT_PRIMARY
    p.font.name = Style.TITLE_FONT

    # Bullet points with square markers
    bullets = config.get("bullets", [])
    start_y = 1.8
    for i, bullet in enumerate(bullets):
        y_pos = start_y + (i * 1.0)

        # Square bullet marker
        marker = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left=Style.MARGIN_LEFT, top=Inches(y_pos + 0.12),
            width=Inches(0.12), height=Inches(0.12)
        )
        marker.fill.solid()
        marker.fill.fore_color.rgb = Style.PRIMARY
        marker.line.fill.background()

        # Bullet text
        text_box = slide.shapes.add_textbox(
            Inches(1.15), Inches(y_pos), Inches(11), Inches(0.8)
        )
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = bullet
        p.font.size = Style.BODY_SIZE
        p.font.color.rgb = Style.TEXT_PRIMARY
        p.font.name = Style.BODY_FONT


def add_closing_slide(prs: Presentation, config: dict):
    """Closing/thank you slide with dark background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, Style.BG_DARK)

    # Title - centered, white
    title_box = slide.shapes.add_textbox(
        Style.MARGIN_LEFT, Inches(2.8), Style.CONTENT_WIDTH, Inches(1.2)
    )
    p = title_box.text_frame.paragraphs[0]
    p.text = config.get("title", "Thank You")
    p.font.size = Pt(52)
    p.font.bold = True
    p.font.color.rgb = Style.TEXT_ON_DARK
    p.font.name = Style.TITLE_FONT
    p.alignment = PP_ALIGN.CENTER

    # Accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Inches(5.5), top=Inches(4.1),
        width=Inches(2.3), height=Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = Style.PRIMARY
    line.line.fill.background()

    # Subtitle (contact info, URL, etc.)
    if "subtitle" in config:
        sub_box = slide.shapes.add_textbox(
            Style.MARGIN_LEFT, Inches(4.5), Style.CONTENT_WIDTH, Inches(0.6)
        )
        p = sub_box.text_frame.paragraphs[0]
        p.text = config["subtitle"]
        p.font.size = Style.SUBTITLE_SIZE
        p.font.color.rgb = Style.TEXT_SECONDARY
        p.font.name = Style.BODY_FONT
        p.alignment = PP_ALIGN.CENTER


def add_blank_slide(prs: Presentation, config: dict):
    """Blank slide for custom content."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, Style.BG_WHITE)

    if "title" in config:
        title_box = slide.shapes.add_textbox(
            Style.MARGIN_LEFT, Inches(0.5), Style.CONTENT_WIDTH, Inches(1)
        )
        p = title_box.text_frame.paragraphs[0]
        p.text = config["title"]
        p.font.size = Style.HEADING_SIZE
        p.font.bold = True
        p.font.color.rgb = Style.PRIMARY


# =============================================================================
# HELPER FUNCTIONS
# =============================================================================

def add_image_to_slide(slide, image_path: str, left: float, top: float,
                       width: float = None, height: float = None):
    """Add an image to a slide at specified position (in inches)."""
    kwargs = {"image_file": image_path, "left": Inches(left), "top": Inches(top)}
    if width:
        kwargs["width"] = Inches(width)
    if height:
        kwargs["height"] = Inches(height)
    slide.shapes.add_picture(**kwargs)


def add_table_to_slide(slide, data: list[list[str]], left: float, top: float,
                       width: float, height: float):
    """Add a table to a slide. First row is styled as header."""
    rows, cols = len(data), len(data[0]) if data else 0
    if rows == 0 or cols == 0:
        return

    table = slide.shapes.add_table(
        rows=rows, cols=cols,
        left=Inches(left), top=Inches(top),
        width=Inches(width), height=Inches(height)
    ).table

    for row_idx, row_data in enumerate(data):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_text)
            if row_idx == 0:
                cell.text_frame.paragraphs[0].font.bold = True


# =============================================================================
# MAIN
# =============================================================================

def main():
    """Generate the presentation and save to file."""
    output_path = Path(sys.argv[1]).expanduser() if len(sys.argv) > 1 else Path(OUTPUT_FILENAME)

    print(f"Creating presentation: {output_path}")
    prs = create_presentation()
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    print(f"Saved to: {output_path.absolute()}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
